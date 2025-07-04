import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

df = pd.read_csv("coffee_chain_cleaned.csv")

ppt = Presentation()

def add_chart_slide(df_plot, title, filename):
    # Save chart image
    plt.tight_layout()
    plt.savefig(filename)
    plt.close()

    slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
    slide.shapes.add_picture(filename, Inches(1), Inches(1.5), width=Inches(8))

df["order_date"] = pd.to_datetime(df["order_date"])
df["year"] = df["order_date"].dt.year
df["month"] = df["order_date"].dt.month

monthly = df.groupby(["year", "month"])["total"].sum().reset_index()


plt.figure(figsize=(8, 5))
cat_rev = df.groupby("category")["total"].sum().sort_values()
sns.barplot(x=cat_rev.values, y=cat_rev.index, palette="viridis")
plt.title("Revenue by Category")
add_chart_slide(cat_rev, "Revenue by Product Category", "chart1.png")

monthly = df.groupby(["year", "month"])["total"].sum().reset_index()
monthly["date"] = pd.to_datetime(monthly[["year", "month"]].assign(day=1))
plt.figure(figsize=(10, 5))
sns.lineplot(x="date", y="total", data=monthly, marker="o")
plt.title("Monthly Revenue Trend")
add_chart_slide(monthly, "Monthly Revenue Trend", "chart2.png")

plt.figure(figsize=(6, 4))
sns.boxplot(data=df, x="gender", y="total", palette="pastel")
plt.title("Order Value by Gender")
add_chart_slide(df, "Order Value by Gender", "chart3.png")

top_products = df.groupby("product_id")["quantity"].sum().nlargest(5).reset_index()
plt.figure(figsize=(8, 5))
sns.barplot(x="quantity", y="product_id", data=top_products, palette="mako")
plt.title("Top 5 Best-Selling Products")
add_chart_slide(top_products, "Top 5 Best-Selling Products", "chart4.png")

region_rev = df.groupby("region")["total"].sum().sort_values()
plt.figure(figsize=(8, 4))
sns.barplot(x=region_rev.values, y=region_rev.index, palette="crest")
plt.title("Revenue by Region")
add_chart_slide(region_rev, "Revenue by Region", "chart5.png")

heat = df.groupby(["region", "category"])["total"].sum().unstack().fillna(0)
plt.figure(figsize=(10, 6))
sns.heatmap(heat, cmap="YlGnBu", annot=True, fmt=".0f")
plt.title("Revenue Heatmap: Region vs Category")
add_chart_slide(heat, "Revenue Heatmap", "chart6.png")

ppt.save("coffee_chain_analysis.pptx")
print("✅ PowerPoint export complete: coffee_chain_analysis.pptx")